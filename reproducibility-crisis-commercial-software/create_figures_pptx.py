#!/usr/bin/env python3
"""Create editable PPTX file with all figures (1 figure per slide, R1 revision)."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = SCRIPT_DIR / "output"
FIG_DIR = OUTPUT_DIR / "figures"

# Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

figures = [
    ('fig1_software_rates_by_field.png',
     'Figure 1',
     'Software mention rates by research field, stratified by software type (N = 10,000 articles, 2020\u20132026).'),
    ('fig2_software_landscape_combined.png',
     'Figure 2',
     '(a) Top 20 software tools declared in published research articles, coloured by licence type. '
     '(b) Usage heatmap across seven research fields with Total row.'),
    ('fig3_version_and_availability.png',
     'Figure 3',
     '(a) Version mention rates among articles reporting software, and (b) code and data availability statement rates, by research field.'),
    ('fig4_version_availability.png',
     'Figure 4',
     'Availability assessment of commercial software versions cited in published articles.'),
    ('fig5_replication_costs.png',
     'Figure 5',
     '(a) Distribution of estimated replication costs among articles using commercial software, and (b) mean replication cost by research field.'),
    ('fig6_pmc_impact.png',
     'Figure 6',
     'Impact of PMC full-text availability on software detection rates.'),
    ('fig7_pmc_subanalysis.png',
     'Figure 7',
     'Sensitivity analysis restricted to the PMC full-text subset.'),
    ('fig8_country_income.png',
     'Figure 8',
     'Software use and reproducibility indicators by country income group (World Bank classification).'),
]

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

for fname, title, caption in figures:
    fig_path = FIG_DIR / fname
    if not fig_path.exists():
        print(f"  WARNING: {fig_path} not found, skipping")
        continue

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Title
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image (centered, scaled)
    img_left = Inches(0.5)
    img_top = Inches(0.9)
    img_width = Inches(12.3)
    img_height = Inches(5.5)
    slide.shapes.add_picture(str(fig_path), img_left, img_top, width=img_width)

    # Caption
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.9))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f'{title}. {caption}'
    p2.font.size = Pt(12)
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER

out_path = OUTPUT_DIR / 'figures_english_r1.pptx'
prs.save(str(out_path))
print(f"Figures PPTX saved: {out_path}")
